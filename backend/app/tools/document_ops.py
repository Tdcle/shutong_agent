"""文档读取工具 — 支持 PDF/DOCX/XLSX/PPTX 等常见文档格式的文字提取。"""

from __future__ import annotations

import json
import os
import subprocess
from pathlib import Path

from app.tools.base import tool
from app.tools.workspace import resolve_readonly

_DOC_SCRIPT = r'''
import json, sys, os, traceback
try:
    file_path = r"{file_path}"
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        import pdfplumber
        with pdfplumber.open(file_path) as pdf:
            pages = []
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    pages.append(f"--- 第{i+1}页 ---\n{text}")
            if not pages:
                print(json.dumps({"error": "PDF 中没有可提取的文字（可能是扫描版图片PDF）。"}))
            else:
                print(json.dumps({"type": "pdf", "pages": len(pages), "text": "\n\n".join(pages)}, ensure_ascii=False))

    elif ext == ".docx":
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        tables_desc = []
        for ti, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            tables_desc.append(f"--- 表格{ti+1} ---\n" + "\n".join(rows))
        text = "\n".join(paragraphs)
        if tables_desc:
            text += "\n\n" + "\n\n".join(tables_desc)
        print(json.dumps({"type": "docx", "paragraphs": len(paragraphs), "tables": len(doc.tables), "text": text}, ensure_ascii=False))

    elif ext in (".xlsx", ".xls"):
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        sheets = []
        for name in wb.sheetnames:
            ws = wb[name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            lines = [f"--- Sheet: {name} ---"]
            for row in rows:
                cells = [str(c) if c is not None else "" for c in row]
                if any(c for c in cells):
                    lines.append(" | ".join(cells))
            sheets.append("\n".join(lines))
        if not sheets:
            print(json.dumps({"error": "Excel 文件中没有可读数据。"}))
        else:
            print(json.dumps({"type": "xlsx", "sheets": len(sheets), "text": "\n\n".join(sheets)}, ensure_ascii=False))

    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for si, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(cells))
            if texts:
                slides.append(f"--- 幻灯片{si+1} ---\n" + "\n".join(texts))
        if not slides:
            print(json.dumps({"error": "PPT 中没有可提取的文字。"}))
        else:
            print(json.dumps({"type": "pptx", "slides": len(slides), "text": "\n\n".join(slides)}, ensure_ascii=False))

    elif ext in (".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml",
                 ".toml", ".ini", ".cfg", ".log", ".html", ".css",
                 ".py", ".js", ".ts", ".java", ".go", ".rs", ".c", ".cpp",
                 ".sh", ".bat", ".sql", ".vue", ".jsx", ".tsx", ".svg"):
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        print(json.dumps({"type": "text", "size": len(text), "text": text}, ensure_ascii=False))

    else:
        print(json.dumps({"error": f"不支持的文件类型: {ext}"}))

except Exception as e:
    print(json.dumps({"error": f"{type(e).__name__}: {e}"}, ensure_ascii=False))
sys.stdout.flush()
'''


def _run_document_reader(file_path: Path) -> str:
    agent_python = str(
        Path(__file__).parent.parent.parent / ".agent_runtime" / "Scripts" / "python.exe"
    )
    if not Path(agent_python).exists():
        return "文档读取环境未找到（.agent_runtime 不存在）。请运行 setup 安装依赖。"

    script = _DOC_SCRIPT.replace("{file_path}", file_path.as_posix())

    clean_env = {}
    strip_prefixes = {
        "PYTHONHOME", "PYTHONPATH", "VIRTUAL_ENV", "CONDA_PREFIX",
        "CONDA_DEFAULT_ENV", "CONDA_PROMPT_MODIFIER",
    }
    for key, val in os.environ.items():
        if key in strip_prefixes:
            continue
        if key.startswith("CONDA_") or key.startswith("PIP_"):
            continue
        clean_env[key] = val
    clean_env.update({
        "PYTHONUNBUFFERED": "1",
        "PYTHONIOENCODING": "utf-8",
    })

    try:
        proc = subprocess.run(
            [agent_python, "-u", "-c", script],
            capture_output=True, text=True, timeout=120,
            encoding="utf-8", errors="replace",
            env=clean_env,
        )
        stdout = proc.stdout.strip()
        stderr = proc.stderr.strip()
    except subprocess.TimeoutExpired:
        return "文档读取超时（超过 120 秒）。文件可能过大。"
    except Exception as exc:
        return f"文档读取异常：{exc}"

    data = None
    if stdout:
        try:
            data = json.loads(stdout)
        except json.JSONDecodeError:
            pass
    if data is None and stderr:
        try:
            data = json.loads(stderr)
        except json.JSONDecodeError:
            pass
    if data is None:
        stderr_preview = stderr[:600] if stderr else "(no output)"
        return f"文档读取失败（无有效输出）。\nstderr: {stderr_preview}"

    if "error" in data:
        return f"文档读取失败：{data['error']}"

    text = data.get("text", "")
    if not text:
        return "文档中没有可提取的文字内容。"

    doc_type = data.get("type", "unknown")
    type_labels = {"pdf": "PDF", "docx": "Word", "xlsx": "Excel", "pptx": "PPT", "text": "文本"}
    label = type_labels.get(doc_type, doc_type.upper())
    pages = data.get("pages") or data.get("slides") or data.get("sheets") or data.get("paragraphs")

    # Truncate very long documents to a reasonable size.
    # For deep analysis, the sub-agent gets the full content via stash (untruncated).
    MAX_CHARS = 50000
    if len(text) > MAX_CHARS:
        text = text[:MAX_CHARS] + f"\n\n... 文档过长，已截断（原文共 {len(text)} 字符，显示前 {MAX_CHARS} 字符）..."

    header = f"[文档内容 — {label}]"
    if pages:
        header += f" 共 {pages} 页/段，{len(text)} 字符"
    footer = "---\n以上是文档的全部文字内容。请直接基于这些信息回答用户，不要再调用其他工具查找资料。"

    return f"{header}\n\n{text}\n\n{footer}"


@tool(
    name="read_document",
    description="读取 PDF、Word、Excel、PPT 等文档的文字内容。结果可直接用于回答用户问题，无需再用 execute_python 二次处理。",
    parameters={
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "文档文件路径"},
        },
        "required": ["path"],
    },
    permission_level="read",
)
def read_document(path: str) -> str:
    p = resolve_readonly(path)
    if not p.exists():
        return f"Error: file does not exist: {path}"
    if not p.is_file():
        return f"Error: path is not a file: {path}"

    ext = p.suffix.lower()
    supported = {
        ".pdf", ".docx", ".xlsx", ".xls", ".pptx",
        ".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml",
        ".toml", ".ini", ".cfg", ".log", ".html", ".css",
        ".py", ".js", ".ts", ".java", ".go", ".rs", ".c", ".cpp",
        ".sh", ".bat", ".sql", ".vue", ".jsx", ".tsx", ".svg",
    }
    if ext not in supported:
        return f"Error: unsupported file type '{ext}'. Supported: PDF, DOCX, XLSX, PPTX, and text-based formats."

    return _run_document_reader(p)
